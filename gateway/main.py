import sys
import os
import uuid
import csv
from datetime import datetime
from pathlib import Path
from io import StringIO
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import List, Dict, Tuple, Optional

sys.path.insert(0, str(Path(__file__).parent.parent))

try:
    import pypdf
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from docx import Document as DocxDocument
    HAS_PYTHON_DOCX = True
except ImportError:
    HAS_PYTHON_DOCX = False

try:
    from pptx import Presentation
    HAS_PYTHON_PPTX = True
except ImportError:
    HAS_PYTHON_PPTX = False

from gateway.ingestion.link_input_handler import LinkInputHandler
from gateway.analysis.hidden_content_analyzer import HiddenContentAnalyzer
from gateway.analysis.prompt_injection_detector import PromptInjectionDetector
from gateway.analysis.exfiltration_detector import ExfiltrationDetector
from gateway.analysis.agentic_intent_detector import AgenticIntentDetector
from gateway.analysis.intent_classifier import IntentClassifier
from gateway.analysis.semantic_detector import SemanticThreatDetector
from gateway.analysis.deobfuscator import ContentDeobfuscator
from gateway.analysis.ocr_analyzer import OCRContentAnalyzer
from gateway.analysis.houyi_pattern_detector import HOUYIPatternDetector
from gateway.analysis.intent_strength_scorer import IntentStrengthScorer, IntentStrength
from gateway.decision_engine.policy_engine import PolicyEngine
from gateway.learning.ml_policy_engine import MLPolicyEngine
from gateway.agent_guard.agent_controller import AgentController
from gateway.shared.schemas import SecurityEvent, SecurityDecision, ContentBlock, RiskLevel
from gateway.shared.logging_utils import SecurityLogger
from gateway.shared.config_loader import get_config

# Objective 3: apply the same ReDoS-aware length discipline at the ingestion
# boundary, before content is handed to any analyzer. This mirrors the caps
# in link_input_handler.py / prompt_injection_detector.py so oversized input
# never even reaches the parsing/regex layers below.
from gateway.shared import safe_regex

# Objective 2: cap on how much of the agent's system prompt we compare
# against. System prompts are typically a few KB; this is generous headroom
# while still bounding SystemPromptContextAnalyzer's tokenization work.
MAX_AGENT_SYSTEM_PROMPT_LENGTH = 20_000

# Optional convenience locations for supplying the protected agent's system
# prompt without changing call sites: an explicit CLI flag always wins, then
# a well-known local file, then an environment variable. All are optional -
# omitting all three preserves prior (context-blind) behavior exactly.
DEFAULT_SYSTEM_PROMPT_FILE = Path("agent_system_prompt.txt")
SYSTEM_PROMPT_ENV_VAR = "PROMPTWALL_AGENT_SYSTEM_PROMPT"


def _resolve_agent_system_prompt(explicit: Optional[str] = None) -> Optional[str]:
    """
    Resolve the protected agent's system prompt from, in priority order:
    1. an explicitly-passed value (e.g. from a caller/API/CLI flag)
    2. a local ./agent_system_prompt.txt file, if present
    3. the PROMPTWALL_AGENT_SYSTEM_PROMPT environment variable
    Returns None if none of the above are available, which is a fully
    supported "no context available" state.
    """
    if explicit:
        return safe_regex.bounded_text(explicit, MAX_AGENT_SYSTEM_PROMPT_LENGTH)

    if DEFAULT_SYSTEM_PROMPT_FILE.exists():
        try:
            content = DEFAULT_SYSTEM_PROMPT_FILE.read_text(encoding="utf-8").strip()
            if content:
                return safe_regex.bounded_text(content, MAX_AGENT_SYSTEM_PROMPT_LENGTH)
        except OSError:
            pass

    env_value = os.environ.get(SYSTEM_PROMPT_ENV_VAR)
    if env_value:
        return safe_regex.bounded_text(env_value, MAX_AGENT_SYSTEM_PROMPT_LENGTH)

    return None


class PromptWall:

    def __init__(self):
        self.config = get_config()

        self.input_handler = LinkInputHandler()
        self.policy_engine = MLPolicyEngine(model_path="/Users/rahulmac/Documents/Projects/projects/Securing_Agentic_AIs/gateway/learning/models/security_model.pkl")
        self.agent_controller = AgentController()
        self.logger = SecurityLogger(log_dir=self.config.get('logging', 'log_dir', default='logs'))

        self.intent_classifier = IntentClassifier()
        self.hidden_analyzer = HiddenContentAnalyzer()
        self.injection_detector = PromptInjectionDetector()
        self.exfiltration_detector = ExfiltrationDetector()
        self.agentic_detector = AgenticIntentDetector()
        self.houyi_detector = HOUYIPatternDetector()

        # Intent strength scorer for better discrimination
        self.intent_strength_scorer = IntentStrengthScorer()

        self.semantic_detector = SemanticThreatDetector() if self.config.is_feature_enabled('semantic_detection') else None
        self.deobfuscator = ContentDeobfuscator() if self.config.is_feature_enabled('active_deobfuscation') else None
        self.ocr_analyzer = OCRContentAnalyzer() if self.config.is_feature_enabled('ocr_analysis') else None

        self.parallel_enabled = self.config.get_parallel_execution_enabled()
        self.max_workers = self.config.get_max_workers()

        self.logger.log_info("Security Wall initialized")
        self.logger.log_info(f"Parallel execution: {self.parallel_enabled}, Max workers: {self.max_workers}")
        self.logger.log_info(f"Semantic detection: {self.semantic_detector is not None}")
        self.logger.log_info(f"De-obfuscation: {self.deobfuscator is not None}")
        self.logger.log_info(f"OCR analysis: {self.ocr_analyzer is not None}")
        # ReDoS mitigation diagnostics (Objective 1): report which regex
        # backend safe_regex actually selected at startup.
        self.logger.log_info(f"Regex backend (ReDoS protection): {safe_regex.active_backend()}")

    def process_input(
        self,
        input_data: str,
        input_type: str = "auto",
        image_path: str = None,
        agent_system_prompt: Optional[str] = None,
    ) -> dict:
        """
        Args:
            input_data: The raw URL or text to evaluate.
            input_type: "auto", "url", or "text".
            image_path: Optional path to an image for OCR analysis.
            agent_system_prompt: Objective 2 - the protected LLM agent's own
                system prompt. When supplied (directly, via
                ./agent_system_prompt.txt, or via the
                PROMPTWALL_AGENT_SYSTEM_PROMPT env var), PolicyEngine weighs
                the user's input against the boundaries and declared purpose
                of that system prompt instead of evaluating it in a vacuum.
                Fully optional; omitting it preserves prior behavior.
        """

        self.logger.log_info(f"Processing input of type: {input_type}")

        resolved_system_prompt = _resolve_agent_system_prompt(agent_system_prompt)
        self.logger.log_info(
            f"Context-aware assessment: "
            f"{'enabled (system prompt supplied)' if resolved_system_prompt else 'disabled (no system prompt supplied)'}"
        )

        # Objective 3: bound the raw input length at the very first
        # opportunity, before it reaches HTML parsing or any regex-based
        # analyzer.
        input_data = safe_regex.bounded_text(input_data or "")

        ocr_text = None
        if image_path and self.ocr_analyzer and self.ocr_analyzer.can_process(image_path):
            self.logger.log_info(f"Performing OCR on image: {image_path}")
            ocr_text = self.ocr_analyzer.extract_text_from_image(image_path)
            if ocr_text:
                self.logger.log_info(f"OCR extracted {len(ocr_text)} characters")
                input_data = input_data + "\n\n[OCR_EXTRACTED_TEXT]\n" + ocr_text

        extracted = self.input_handler.process_input(input_data, input_type)

        decoded_content = ""
        if self.deobfuscator:
            decoded_content = self.deobfuscator.get_decoded_content(
                extracted.visible_text,
                extracted.hidden_elements
            )
            if decoded_content:
                self.logger.log_info(f"De-obfuscation extracted {len(decoded_content)} characters")
                extracted.hidden_elements.append(f"[DECODED_CONTENT] {decoded_content}")

        content_blocks = [
            ContentBlock(
                content=extracted.visible_text,
                content_type="text",
                visibility="visible"
            )
        ]

        for idx, hidden in enumerate(extracted.hidden_elements):
            content_blocks.append(
                ContentBlock(
                    content=hidden,
                    content_type="html",
                    visibility="hidden",
                    source_location=f"hidden_element_{idx}"
                )
            )

        if self.parallel_enabled:
            analysis_results = self._run_parallel_analysis(
                extracted.visible_text,
                extracted.hidden_elements,
                extracted.metadata,
                ocr_text,
                extracted.location_map
            )
        else:
            analysis_results = self._run_sequential_analysis(
                extracted.visible_text,
                extracted.hidden_elements,
                extracted.metadata,
                ocr_text,
                extracted.location_map
            )

        # Objective 2: pass agent_system_prompt through to the decision
        # engine. MLPolicyEngine (a PolicyEngine subclass) may not yet
        # accept this kwarg if it hasn't been updated in lockstep with
        # policy_engine.py, so we degrade gracefully rather than crash the
        # whole request - context-aware assessment simply won't apply for
        # that call, and this is logged so it's not silently invisible.
        try:
            assessment = self.policy_engine.make_decision(
                analysis_results,
                extracted.visible_text,
                extracted.hidden_elements,
                raw_input=input_data,
                agent_system_prompt=resolved_system_prompt,
            )
        except TypeError:
            self.logger.log_error(
                "PolicyEngine.make_decision() does not accept 'agent_system_prompt' "
                "yet on this build; retrying without context-aware assessment."
            )
            assessment = self.policy_engine.make_decision(
                analysis_results,
                extracted.visible_text,
                extracted.hidden_elements,
                raw_input=input_data,
            )

        assessment.content_blocks = content_blocks
        assessment.source = extracted.metadata.get("source_url", "direct_input")

        agentic_analysis = next(
            (r for r in analysis_results if r.module_name == "agentic_intent_detector"),
            None
        )

        if agentic_analysis and agentic_analysis.risk_level not in [RiskLevel.SAFE, RiskLevel.LOW]:
            assessment.agentic_intent_detected = True
            requested_actions = [
                f.get('action') for f in agentic_analysis.findings
                if f.get('type') == 'action_request' and f.get('action')
            ]
            assessment.requested_actions = list(set(requested_actions))

        session_id = str(uuid.uuid4())

        restrictions = self.policy_engine._determine_restrictions(
            assessment.overall_risk,
            analysis_results,
            assessment.primary_intent
        )

        apply_result = self.agent_controller.apply_restrictions(session_id, restrictions)

        self.logger.log_info(
            f"Restrictions applied: mode={restrictions.mode}, enforcement={apply_result.get('enforcement_status')}"
        )

        event = SecurityEvent(
            event_id=assessment.input_id,
            timestamp=assessment.timestamp.isoformat(),
            event_type="input_processed",
            severity=assessment.overall_risk.value,
            input_source=assessment.source,
            risk_level=assessment.overall_risk.value,
            decision=assessment.decision.value,
            findings=[
                {
                    "module": r.module_name,
                    "risk": r.risk_level.value,
                    "confidence": r.confidence,
                    "details": r.details,
                    "risk_score": r.risk_score,
                    "detected_intent": r.detected_intent.value if r.detected_intent else None
                }
                for r in analysis_results
            ],
            metadata={
                "session_id": session_id,
                "risk_score": assessment.risk_score,
                "restricted_capabilities": assessment.restricted_capabilities,
                "agentic_intent_detected": assessment.agentic_intent_detected,
                "requested_actions": assessment.requested_actions,
                "enforcement_mode": restrictions.mode,
                "requires_approval": restrictions.requires_approval,
                "primary_intent": assessment.primary_intent.value,
                "intent_confidence": assessment.intent_confidence,
                "ocr_used": ocr_text is not None,
                "decoded_content_found": len(decoded_content) > 0 if decoded_content else False,
                "context_aware_assessment_used": resolved_system_prompt is not None,
            }
        )

        self.logger.log_security_event(event)

        return self._format_response(assessment, session_id)

    def _run_parallel_analysis(
        self,
        visible_text: str,
        hidden_elements: List[str],
        metadata: dict,
        ocr_text: str = None,
        location_map: dict = None
    ) -> List:
        """
        Run analysis in parallel with intent-aware hidden content analysis.
        """

        analysis_results = []

        # First, run intent classification sequentially to get intent strength
        intent_analysis = self.intent_classifier.analyze(visible_text, hidden_elements)
        analysis_results.append(intent_analysis)

        # Determine intent strength for hidden content analysis
        intent_strength = IntentStrength.WEAK  # Default
        if intent_analysis.detected_intent:
            intent_strength, _ = self.intent_strength_scorer.score_intent_strength(
                visible_text,
                intent_analysis.detected_intent.value
            )

        # Now run other analyses in parallel
        tasks = [
            ("hidden", lambda: self.hidden_analyzer.analyze(
                visible_text, hidden_elements, location_map, intent_strength
            )),
            ("injection", lambda: self.injection_detector.analyze(visible_text, hidden_elements, location_map)),
            ("exfiltration", lambda: self.exfiltration_detector.analyze(visible_text, hidden_elements, metadata)),
            ("agentic", lambda: self.agentic_detector.analyze(visible_text, hidden_elements)),
            ("houyi", lambda: self.houyi_detector.analyze(visible_text, hidden_elements)),
        ]

        if self.semantic_detector:
            tasks.append(("semantic", lambda: self.semantic_detector.analyze(visible_text, hidden_elements)))

        if self.deobfuscator:
            tasks.append(("deobfuscator", lambda: self.deobfuscator.analyze(visible_text, hidden_elements)))

        if self.ocr_analyzer and ocr_text:
            tasks.append(("ocr", lambda: self.ocr_analyzer.analyze(ocr_text, "image_input")))

        with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
            future_to_task = {executor.submit(task_fn): task_name for task_name, task_fn in tasks}

            for future in as_completed(future_to_task):
                task_name = future_to_task[future]
                try:
                    result = future.result(timeout=self.config.get('performance', 'analysis_timeout', default=30))
                    analysis_results.append(result)
                    self.logger.log_info(f"Completed analysis: {task_name}")
                except Exception as e:
                    self.logger.log_error(f"Analysis failed for {task_name}: {str(e)}")

        return analysis_results

    def _run_sequential_analysis(
        self,
        visible_text: str,
        hidden_elements: List[str],
        metadata: dict,
        ocr_text: str = None,
        location_map: dict = None
    ) -> List:
        """
        Run analysis sequentially with intent-aware hidden content analysis.
        """

        analysis_results = []

        # First, determine intent and its strength
        intent_analysis = self.intent_classifier.analyze(visible_text, hidden_elements)
        analysis_results.append(intent_analysis)

        intent_strength = IntentStrength.WEAK  # Default
        if intent_analysis.detected_intent:
            intent_strength, _ = self.intent_strength_scorer.score_intent_strength(
                visible_text,
                intent_analysis.detected_intent.value
            )

        # Run hidden analysis with intent strength awareness
        hidden_analysis = self.hidden_analyzer.analyze(visible_text, hidden_elements, location_map, intent_strength)
        analysis_results.append(hidden_analysis)

        # Continue with other analyses
        injection_analysis = self.injection_detector.analyze(visible_text, hidden_elements, location_map)
        analysis_results.append(injection_analysis)

        exfiltration_analysis = self.exfiltration_detector.analyze(visible_text, hidden_elements, metadata)
        analysis_results.append(exfiltration_analysis)

        agentic_analysis = self.agentic_detector.analyze(visible_text, hidden_elements)
        analysis_results.append(agentic_analysis)

        houyi_analysis = self.houyi_detector.analyze(visible_text, hidden_elements)
        analysis_results.append(houyi_analysis)

        if self.semantic_detector:
            semantic_analysis = self.semantic_detector.analyze(visible_text, hidden_elements)
            analysis_results.append(semantic_analysis)

        if self.deobfuscator:
            deobfuscator_analysis = self.deobfuscator.analyze(visible_text, hidden_elements)
            analysis_results.append(deobfuscator_analysis)

        if self.ocr_analyzer and ocr_text:
            ocr_analysis = self.ocr_analyzer.analyze(ocr_text, "image_input")
            analysis_results.append(ocr_analysis)

        return analysis_results

    def _format_response(self, assessment, session_id: str) -> dict:

        return {
            "session_id": session_id,
            "input_id": assessment.input_id,
            "timestamp": assessment.timestamp.isoformat(),
            "decision": assessment.decision.value,
            "risk_level": assessment.overall_risk.value,
            "risk_score": round(assessment.risk_score, 3),
            "primary_intent": assessment.primary_intent.value,
            "intent_confidence": round(assessment.intent_confidence, 3),
            "agentic_intent_detected": assessment.agentic_intent_detected,
            "requested_actions": assessment.requested_actions,
            "content": {
                "original": assessment.content_blocks[0].content if assessment.content_blocks else "",
                "sanitized": assessment.sanitized_content,
                "hidden_elements_count": len([b for b in assessment.content_blocks if b.visibility == "hidden"])
            },
            "analysis": [
                {
                    "module": r.module_name,
                    "risk": r.risk_level.value,
                    "confidence": round(r.confidence, 3),
                    "findings_count": len(r.findings),
                    "details": r.details,
                    "risk_score": round(r.risk_score, 3),
                    "detected_intent": r.detected_intent.value if r.detected_intent else None
                }
                for r in assessment.analysis_results
            ],
            "restrictions": assessment.restricted_capabilities,
            "reasoning": assessment.reasoning,
            "allowed": assessment.decision == SecurityDecision.ALLOW
        }


def read_file_content(file_path: Path) -> Tuple[str, str]:

    suffix = file_path.suffix.lower()

    image_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp'}
    if suffix in image_extensions:
        return f"[IMAGE FILE: {file_path.name}]", str(file_path)

    text_extensions = {'.txt', '.md', '.html', '.xml', '.csv', '.json', '.py', '.js', '.yaml', '.yml', '.log'}
    if suffix in text_extensions:
        try:
            return file_path.read_text(encoding='utf-8').strip(), None
        except UnicodeDecodeError:
            try:
                return file_path.read_text(encoding='latin-1').strip(), None
            except Exception as e:
                raise ValueError(f"Cannot read text file {file_path}: {str(e)}")

    elif suffix == '.csv':
        try:
            content = []
            with open(file_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                for row_num, row in enumerate(reader, 1):
                    content.append(f"Row {row_num}: {' | '.join(row)}")
            result = '\n'.join(content).strip()
            if not result:
                raise ValueError("CSV file is empty")
            return result, None
        except Exception as e:
            raise ValueError(f"Error reading CSV file {file_path}: {str(e)}")

    elif suffix in {'.xlsx', '.xls'}:
        if not HAS_OPENPYXL:
            raise ValueError(
                f"Excel file detected ({file_path}) but openpyxl is not installed. "
                "Install it with: pip install openpyxl"
            )
        try:
            content = []
            workbook = openpyxl.load_workbook(file_path)
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content.append(f"Sheet: {sheet_name}")
                for row_num, row in enumerate(sheet.iter_rows(values_only=True), 1):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    content.append(f"Row {row_num}: {' | '.join(row_values)}")
            result = '\n'.join(content).strip()
            if not result:
                raise ValueError("Excel file is empty or contains no readable data")
            return result, None
        except Exception as e:
            raise ValueError(f"Error reading Excel file {file_path}: {str(e)}")

    elif suffix == '.docx':
        if not HAS_PYTHON_DOCX:
            raise ValueError(
                f"Word document detected ({file_path}) but python-docx is not installed. "
                "Install it with: pip install python-docx"
            )
        try:
            doc = DocxDocument(file_path)
            content = []
            for para in doc.paragraphs:
                if para.text.strip():
                    content.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    row_values = [cell.text for cell in row.cells]
                    content.append(' | '.join(row_values))
            result = '\n'.join(content).strip()
            if not result:
                raise ValueError("Word document contains no readable text")
            return result, None
        except Exception as e:
            raise ValueError(f"Error reading Word document {file_path}: {str(e)}")

    elif suffix == '.pptx':
        if not HAS_PYTHON_PPTX:
            raise ValueError(
                f"PowerPoint file detected ({file_path}) but python-pptx is not installed. "
                "Install it with: pip install python-pptx"
            )
        try:
            prs = Presentation(file_path)
            content = []
            for slide_num, slide in enumerate(prs.slides, 1):
                content.append(f"Slide {slide_num}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content.append(f"  {shape.text}")
            result = '\n'.join(content).strip()
            if not result:
                raise ValueError("PowerPoint presentation contains no readable text")
            return result, None
        except Exception as e:
            raise ValueError(f"Error reading PowerPoint file {file_path}: {str(e)}")

    elif suffix == '.pdf':
        if not HAS_PYPDF:
            raise ValueError(
                f"PDF file detected ({file_path}) but pypdf is not installed. "
                "Install it with: pip install pypdf"
            )
        try:
            text_content = []
            with open(file_path, 'rb') as f:
                pdf_reader = pypdf.PdfReader(f)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text_content.append(page.extract_text())
            result = '\n'.join(text_content).strip()
            if not result:
                raise ValueError("No text could be extracted from the PDF")
            return result, None
        except Exception as e:
            raise ValueError(f"Error reading PDF file {file_path}: {str(e)}")

    elif suffix in {'.zip', '.tar', '.gz', '.bin'}:
        raise ValueError(
            f"Binary file format '{suffix}' is not supported. "
            "Supported: text, CSV, Excel, Word, PowerPoint, PDF, and images for OCR."
        )

    else:
        try:
            return file_path.read_text(encoding='utf-8').strip(), None
        except (UnicodeDecodeError, OSError):
            raise ValueError(
                f"Unable to read file {file_path}. "
                "Supported: .txt, .md, .html, .csv, .xlsx, .docx, .pptx, .pdf, and images (.jpg, .png, etc.)"
            )


def main():

    print("=" * 60)
    print("PromptWall - Enhanced LLM Security Gateway")
    print("Features: HOUYI Detection | Semantic | De-obfuscation | OCR | ReDoS-Safe Regex | Context-Aware Policy")
    print("=" * 60)
    print()

    guard = PromptWall()

    # Objective 2: optional CLI flag to point at the protected agent's
    # system prompt file, e.g.:
    #   python main.py --system-prompt agent_prompt.txt "some input text"
    # This is resolved in addition to (and takes priority over) the
    # ./agent_system_prompt.txt / env var fallbacks in
    # _resolve_agent_system_prompt().
    argv = sys.argv[1:]
    cli_system_prompt = None
    if "--system-prompt" in argv:
        idx = argv.index("--system-prompt")
        if idx + 1 < len(argv):
            sp_path = Path(argv[idx + 1])
            if sp_path.exists() and sp_path.is_file():
                try:
                    cli_system_prompt = sp_path.read_text(encoding='utf-8').strip()
                except OSError as e:
                    print(f"Warning: could not read --system-prompt file: {e}", file=sys.stderr)
            # Remove the flag and its value from argv so it isn't treated as input
            del argv[idx:idx + 2]

    if argv:
        input_arg = " ".join(argv)

        file_path = Path(input_arg)
        if file_path.exists() and file_path.is_file():
            try:
                file_content, image_path = read_file_content(file_path)
                if file_content:
                    print(f"Processing input from file: {file_path}")
                    if image_path:
                        print(f"OCR will be performed on image")
                    print("-" * 60 + "\n")
                    process_and_display_result(guard, file_content, str(file_path), image_path, cli_system_prompt)
                    return
            except ValueError as e:
                print(f"Error: {str(e)}", file=sys.stderr)
                sys.exit(1)
        else:
            process_and_display_result(guard, input_arg, agent_system_prompt=cli_system_prompt)
            return

    if not sys.stdin.isatty():
        stdin_content = sys.stdin.read().strip()

        if not stdin_content:
            print("Error: No input provided via stdin", file=sys.stderr)
            sys.exit(1)

        print("Processing input from stdin...")
        print("-" * 60 + "\n")

        process_and_display_result(guard, stdin_content, agent_system_prompt=cli_system_prompt)
        return

    input_file_path = Path("input.txt")
    if input_file_path.exists():
        try:
            file_content, image_path = read_file_content(input_file_path)
            if file_content:
                print("Processing input from input.txt...")
                print("-" * 60 + "\n")
                process_and_display_result(guard, file_content, "input.txt", image_path, cli_system_prompt)
                return
        except ValueError as e:
            print(f"Error: {str(e)}", file=sys.stderr)
            sys.exit(1)

    print("Enter URL, text, or file path (or 'quit' to exit):")
    print("Supported: .txt, .md, .html, .csv, .xlsx, .docx, .pptx, .pdf, images (.jpg, .png)")
    print()

    while True:
        try:
            user_input = input("> ").strip()

            if user_input.lower() in ['quit', 'exit', 'q']:
                print("\nExiting PromptWall. Stay secure!")
                break

            if not user_input:
                continue

            file_path = Path(user_input)
            if file_path.exists() and file_path.is_file():
                try:
                    file_content, image_path = read_file_content(file_path)
                    if file_content:
                        print("\n" + "-" * 60)
                        print(f"Processing file: {file_path}")
                        if image_path:
                            print("OCR will be performed on image")
                        print("-" * 60 + "\n")
                        result = guard.process_input(file_content, image_path=image_path, agent_system_prompt=cli_system_prompt)
                        display_result(result)
                    else:
                        print(f"File {file_path} is empty.")
                except ValueError as e:
                    print(f"Error: {str(e)}")
                continue

            print("\n" + "-" * 60)
            print("Processing input...")
            print("-" * 60 + "\n")

            result = guard.process_input(user_input, agent_system_prompt=cli_system_prompt)
            display_result(result)

        except KeyboardInterrupt:
            print("\n\nExiting PromptWall. Stay secure!")
            break
        except Exception as e:
            print(f"\nError: {str(e)}")
            print()


def process_and_display_result(guard, input_data, source=None, image_path=None, agent_system_prompt=None):
    try:
        if source:
            print(f"File: {source}")
        result = guard.process_input(input_data, image_path=image_path, agent_system_prompt=agent_system_prompt)
        display_result(result)
    except Exception as e:
        print(f"Error: {str(e)}", file=sys.stderr)
        sys.exit(1)


def display_result(result):
    print(f"Decision: {result['decision'].upper()}")
    print(f"Risk Level: {result['risk_level'].upper()}")
    print(f"Risk Score: {result['risk_score']}")
    print(f"Primary Intent: {result['primary_intent'].upper()} (confidence: {result['intent_confidence']})")

    if result['agentic_intent_detected']:
        print(f"\nAGENTIC INTENT DETECTED")
        if result['requested_actions']:
            print(f"Requested Actions: {', '.join(result['requested_actions'])}")

    print()

    if result['analysis']:
        print("Analysis Results:")
        for analysis in result['analysis']:
            print(f"  - {analysis['module']}: {analysis['risk']} "
                  f"(confidence: {analysis['confidence']}, "
                  f"findings: {analysis['findings_count']}, "
                  f"risk_score: {analysis['risk_score']})")
        print()

    if result['restrictions']:
        print("Restrictions Applied:")
        for restriction in result['restrictions']:
            print(f"  - {restriction}")
        print()

    print("Reasoning:")
    print(result['reasoning'])
    print()

    if result['content']['sanitized']:
        sanitized = result['content']['sanitized']
        print("Sanitized Content:")
        print(sanitized[:200] + "..." if len(sanitized) > 200 else sanitized)
        print()

    print("=" * 60)
    print()


if __name__ == "__main__":
    main()